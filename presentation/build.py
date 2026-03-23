"""
build.py — Generates WasteHero_Driver_Auth_POC.pptx
Rule: every textbox height = (number_of_lines × 0.35") + 0.15" padding
      no element may start within 0.1" of another element's bottom
Run: python3 build.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── Colors ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0f, 0x16, 0x24)
CARD     = RGBColor(0x1a, 0x25, 0x35)
CARD_ALT = RGBColor(0x12, 0x1c, 0x2c)
ACCENT   = RGBColor(0x00, 0xc9, 0xd7)
BODY     = RGBColor(0xc9, 0xd4, 0xe8)
MUTED    = RGBColor(0x8a, 0x96, 0xaa)
WHITE    = RGBColor(0xff, 0xff, 0xff)
DARK     = RGBColor(0x0a, 0x11, 0x1e)
BORDER   = RGBColor(0x25, 0x2f, 0x42)
GREEN_BG = RGBColor(0x0d, 0x28, 0x1a)
GREEN_FG = RGBColor(0x48, 0xbb, 0x78)
RED_BG   = RGBColor(0x2a, 0x10, 0x10)
RED_FG   = RGBColor(0xff, 0x6b, 0x6b)

LOGO = "/Users/mohamedelsdody/Desktop/WastHero/wastehero_mobileapp_navigator/src/images/logo/logo.png"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(BLANK)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl


def box(sl, x, y, w, h, fill=CARD, border=None):
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(0.8)
    else:
        sh.line.fill.background()
    return sh


def txt(sl, text, x, y, w, h,
        size=16, bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False          # no wrap → no overflow surprises
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name  = "Calibri"
    return tb


def txt_wrap(sl, text, x, y, w, h,
             size=14, bold=False, color=BODY, align=PP_ALIGN.LEFT):
    """Textbox with word-wrap enabled. h must be generous enough to hold all lines."""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.name  = "Calibri"
    return tb


def bullets(sl, lines, x, y, w, size=15, color=BODY, gap=0.42):
    """
    Renders each line as its own single-line textbox stacked top-down.
    Returns the y position after the last item.
    """
    cy = y
    for line in lines:
        txt(sl, line, x, cy, w, gap, size=size, color=color)
        cy += gap
    return cy


def top_bar(sl):
    b = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.07))
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background()


def slide_header(sl, title, sub=None):
    top_bar(sl)
    txt(sl, title, 0.45, 0.14, 12.4, 0.62, size=30, bold=True, color=ACCENT)
    if sub:
        txt(sl, sub, 0.45, 0.78, 12.4, 0.30, size=13, color=MUTED, italic=True)


def footer(sl):
    box(sl, 0, 7.18, 13.33, 0.32, fill=DARK)
    txt(sl, "WasteHero Fleet   |   Driver Authentication POC   |   March 2026",
        0.45, 7.2, 12.4, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER, italic=True)


def divider(sl, y, color=BORDER):
    b = sl.shapes.add_shape(1, Inches(0.45), Inches(y), Inches(12.43), Inches(0.02))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()


def numbered_card(sl, num_str, title, body_lines, x, y, w,
                  num_color=ACCENT, card_h=None, title_size=15, body_size=13):
    """
    Horizontal card with left number badge.
    card_h: explicit height override. If None, auto-calculated.
    Returns bottom y of card.
    """
    LINE_H  = 0.32
    title_h = 0.36
    auto_h  = title_h + 0.12 + len(body_lines) * LINE_H + 0.10 + 0.16
    ch      = card_h if card_h else auto_h

    box(sl, x, y, w, ch, fill=CARD, border=BORDER)
    box(sl, x, y, 0.50, ch, fill=num_color)
    txt(sl, num_str, x, y + (ch - 0.34) / 2, 0.50, 0.34,
        size=15, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, title, x + 0.64, y + 0.10, w - 0.76, title_h,
        size=title_size, bold=True, color=WHITE)
    cy = y + title_h + 0.20
    for line in body_lines:
        txt(sl, line, x + 0.64, cy, w - 0.76, LINE_H + 0.04, size=body_size, color=BODY)
        cy += LINE_H + 0.04
    return y + ch


def step_list(sl, steps, x, y, w):
    """
    Vertical step list: (num, title, detail) per step.
    Returns final y.
    """
    LINE_H = 0.36
    cy = y
    for num, title, detail in steps:
        # number circle placeholder
        box(sl, x, cy, 0.45, 0.38, fill=ACCENT)
        txt(sl, num, x, cy, 0.45, 0.38,
            size=15, bold=True, color=BG, align=PP_ALIGN.CENTER)
        txt(sl, title, x + 0.58, cy, w - 0.7, 0.38,
            size=15, bold=True, color=WHITE)
        if detail:
            txt(sl, detail, x + 0.58, cy + 0.38, w - 0.7, 0.30,
                size=12, color=MUTED)
            cy += 0.38 + 0.30 + 0.18
        else:
            cy += 0.38 + 0.15
    return cy


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()

# Full-height left panel
box(s, 0, 0, 4.7, 7.5, fill=CARD)
# Thin top accent bar
box(s, 0, 0, 13.33, 0.07, fill=ACCENT)

# Logo
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.5), Inches(0.3), Inches(3.1), Inches(1.1))

# Left panel: author block
divider(s, 1.65)
txt(s, "Prepared by",       0.5, 1.8,  3.7, 0.32, size=12, color=MUTED)
txt(s, "Mohamed Elsdody",   0.5, 2.15, 3.7, 0.52, size=20, bold=True, color=WHITE)
txt(s, "Mobile Architect",  0.5, 2.68, 3.7, 0.36, size=15, color=ACCENT)
divider(s, 3.18)
txt(s, "March 2026",        0.5, 3.32, 3.7, 0.32, size=13, color=MUTED)
txt(s, "CONFIDENTIAL",      0.5, 3.72, 3.7, 0.30, size=11, color=MUTED, italic=True)

# Right side: title block
txt(s, "Driver Authentication",         5.1, 1.8,  7.8, 0.92, size=44, bold=True,  color=WHITE)
txt(s, "POC Results & Architecture Review", 5.1, 2.8,  7.8, 0.52, size=20, color=ACCENT)
divider(s, 3.48)
txt(s, "Two solutions evaluated against WebAuthn passkeys", 5.1, 3.62, 7.8, 0.38, size=15, color=MUTED)
txt(s, "Recommendation included",       5.1, 4.04, 7.8, 0.34, size=14, color=MUTED, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Agenda")

items = [
    ("01", "The Problem"),
    ("02", "Why Passkey (WebAuthn) Was Rejected"),
    ("03", "Solution A  —  NFC Card + PIN"),
    ("04", "Solution B  —  PIN + Face Verification"),
    ("05", "Head-to-Head Comparison"),
    ("06", "Recommendation & Next Steps"),
]

ITEM_H = 0.72
START_Y = 1.2
for i, (num, label) in enumerate(items):
    y = START_Y + i * ITEM_H
    box(s, 0.45, y, 0.95, 0.55, fill=ACCENT)
    txt(s, num, 0.45, y, 0.95, 0.55, size=17, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, 1.55, y, 11.33, 0.55, fill=CARD, border=BORDER)
    txt(s, label, 1.75, y + 0.09, 11.0, 0.38, size=17, color=BODY)

footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Problem
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "The Problem", "What drove this proof of concept")

# 3 cards with generous spacing
problems = [
    ("01", "Shared Tablets Across Drivers",
     ["Multiple drivers share the same tablets across shifts and depots.",
      "Any enrolled driver must log in on any tablet — no device-specific credentials."]),
    ("02", "Field Workers Need Speed",
     ["Drivers work in the field, often with gloves or in adverse conditions.",
      "Authentication must complete in under 10 seconds. Passwords are impractical."]),
    ("03", "Identity Verification Required",
     ["A PIN alone is not enough — buddy-punching must be prevented.",
      "The system must confirm the person entering the PIN is the registered driver."]),
]

cy = 1.2
for num, title, lines in problems:
    cy = numbered_card(s, num, title, lines, 0.45, cy, 12.43) + 0.22
footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why Not Passkey
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Why Not Passkey (WebAuthn)?", "Evaluated and ruled out early in the project")

# Verdict banner
box(s, 0.45, 1.2, 12.43, 0.52, fill=RED_BG, border=RED_FG)
txt(s, "WebAuthn / FIDO2 Passkeys do not fit the fleet authentication model",
    0.65, 1.27, 12.1, 0.38, size=15, bold=True, color=RED_FG)

reasons = [
    ("01", "Hardware Dependency",
     ["Requires a FIDO2-certified platform authenticator (fingerprint chip / secure enclave).",
      "Fleet ruggedised Android tablets may not carry FIDO2-compliant hardware."]),
    ("02", "Device-Bound by Design",
     ["Passkeys are cryptographically tied to the device they were created on.",
      "Driver registered on Tablet A cannot authenticate on Tablet B — breaks fleet model."]),
    ("03", "No Fleet Management Model",
     ["Designed for consumer web login: one user, one device, one site.",
      "No admin enrollment, no driver revocation, no fleet-wide credential distribution."]),
    ("04", "Wrong Problem Domain",
     ["WebAuthn solves phishing-resistant browser login for individual users.",
      "Fleet driver identity verification is a different problem requiring a different tool."]),
]

# 4 cards: available height = 7.18 - 1.88 = 5.30". Each card = 1.20", gap = 0.12"
cy = 1.88
for num, title, lines in reasons:
    cy = numbered_card(s, num, title, lines, 0.45, cy, 12.43,
                       num_color=MUTED, card_h=1.20) + 0.12
footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Solution A: NFC Overview
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Solution A — NFC Card + PIN", "Branch: feature/nfc-auth")

# Tag
box(s, 0.45, 1.18, 2.4, 0.38, fill=RGBColor(0x10, 0x38, 0x38), border=ACCENT)
txt(s, "IMPLEMENTED", 0.45, 1.18, 2.4, 0.38, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

steps = [
    ("1", "Admin Issues Card",
     "Admin writes driver ID to NFC card via AdminScreen (NDEF text record)."),
    ("2", "Driver Taps Card",
     "Driver holds NFC card near tablet. App reads driver ID from NDEF record."),
    ("3", "PIN Entry",
     "Driver enters PIN. App posts SHA-256 hash to cloud API for verification."),
    ("4", "Access Granted",
     "On match, driver is routed to Dashboard. Session timestamp is logged."),
]

# Horizontal 4-box flow with generous widths
BOX_W = 2.82
BOX_H = 2.95
GAP   = 0.38
SX    = 0.45
SY    = 1.72

for i, (num, title, detail) in enumerate(steps):
    x = SX + i * (BOX_W + GAP)
    box(s, x, SY, BOX_W, BOX_H, fill=CARD, border=BORDER)
    # header strip
    box(s, x, SY, BOX_W, 0.48, fill=ACCENT)
    txt(s, num, x, SY, BOX_W, 0.48,
        size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.12, SY + 0.58, BOX_W - 0.24, 0.40,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt_wrap(s, detail, x + 0.12, SY + 1.08, BOX_W - 0.24, 1.72,
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    # arrow
    if i < 3:
        ax = x + BOX_W + 0.12
        txt(s, ">", ax, SY + 1.04, 0.14, 0.38,
            size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# Tech strip — well below boxes
box(s, 0.45, 4.88, 12.43, 0.46, fill=DARK, border=BORDER)
txt(s, "Tech:  Android NFC API   |   NDEF Text Records   |   Neon PostgreSQL   |   Vercel",
    0.65, 4.94, 12.1, 0.34, size=13, color=MUTED)
footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Solution A: Pros & Cons
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Solution A — Strengths & Limitations")

# Left column: Strengths
box(s, 0.45, 1.18, 5.9, 0.46, fill=GREEN_BG, border=GREEN_FG)
txt(s, "STRENGTHS", 0.45, 1.18, 5.9, 0.46, size=14, bold=True, color=GREEN_FG, align=PP_ALIGN.CENTER)

pros = [
    "+ Fast UX — card tap takes under 1 second",
    "+ Physical 2nd factor (something you have + something you know)",
    "+ Familiar to drivers — similar to existing access-card systems",
    "+ Low implementation complexity — simple NDEF + REST API",
    "+ No biometric hardware required on the tablet",
]

cy = 1.78
for line in pros:
    box(s, 0.45, cy, 5.9, 0.42, fill=CARD, border=BORDER)
    txt(s, line, 0.62, cy + 0.04, 5.6, 0.34, size=13, color=BODY)
    cy += 0.46

# Right column: Limitations
box(s, 6.98, 1.18, 5.9, 0.46, fill=RED_BG, border=RED_FG)
txt(s, "LIMITATIONS", 6.98, 1.18, 5.9, 0.46, size=14, bold=True, color=RED_FG, align=PP_ALIGN.CENTER)

cons = [
    "- Card can be lost, stolen, or cloned — no crypto binding",
    "- PIN stored as plain text in DB (POC only — must be fixed)",
    "- Requires card printing, distribution and re-issuance workflow",
    "- No liveness check — anyone with the card can attempt PIN",
    "- No card revocation mechanism in current implementation",
]

cy = 1.78
for line in cons:
    box(s, 6.98, cy, 5.9, 0.42, fill=CARD_ALT, border=BORDER)
    txt(s, line, 7.15, cy + 0.04, 5.6, 0.34, size=13, color=BODY)
    cy += 0.46

footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Solution B: PIN + Face Overview
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Solution B — PIN + Face Verification", "Branch: feature/facenet-auth   |   Recommended")

box(s, 0.45, 1.18, 2.8, 0.38, fill=GREEN_BG, border=GREEN_FG)
txt(s, "RECOMMENDED", 0.45, 1.18, 2.8, 0.38, size=12, bold=True, color=GREEN_FG, align=PP_ALIGN.CENTER)

# 5 steps as a vertical list on left, detail text on right
steps = [
    ("1", "PIN Entry",
     "Driver enters 6-digit PIN.  SHA-256 + pepper hash is sent to cloud."),
    ("2", "Cloud Driver Lookup",
     "POST /api/fleet/verify { pinHash }  →  returns driverId + name, or 404."),
    ("3", "Face Capture (On-Device)",
     "FaceNet TFLite runs locally.  Blink liveness check.  512-float embedding produced."),
    ("4", "Cloud Face Verification",
     "POST /api/fleet/verify { pinHash, embedding }  →  server decrypts stored vectors + cosine similarity."),
    ("5", "Access Decision",
     "Score >= 0.75  →  Dashboard.   Fail  →  lockout counter incremented (5 fails = 60 min lock)."),
]

SY = 1.72
ROW_H = 0.96
for i, (num, title, detail) in enumerate(steps):
    y = SY + i * ROW_H
    box(s, 0.45, y, 12.43, ROW_H - 0.1, fill=CARD, border=BORDER)
    box(s, 0.45, y, 0.48, ROW_H - 0.1, fill=ACCENT)
    txt(s, num, 0.45, y + 0.17, 0.48, 0.40,
        size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, title, 1.08, y + 0.08, 3.8, 0.36, size=14, bold=True, color=WHITE)
    txt(s, detail, 5.0, y + 0.08, 7.7, 0.70, size=12, color=MUTED)
    if i < 4:
        divider(s, y + ROW_H - 0.1)

# Tech strip
box(s, 0.45, 6.6, 12.43, 0.40, fill=DARK, border=BORDER)
txt(s, "Tech:  TFLite FaceNet   |   ML Kit   |   AES-256-GCM   |   OkHttp3   |   Neon PostgreSQL   |   Vercel",
    0.65, 6.65, 12.1, 0.30, size=12, color=MUTED)
footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — How Face Detection Works
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "How Face Detection & Verification Works",
             "On-device ML pipeline — no face images leave the device")

pipeline = [
    ("1", "Camera Frame",      "CameraX streams live frames to the ML pipeline continuously."),
    ("2", "ML Kit Detection",  "Face bounding box extracted. Eye open probability measured per frame."),
    ("3", "Blink Liveness",    "Eyes open  ->  closed  ->  open sequence required before accepting frame."),
    ("4", "FaceNet TFLite",    "Face region cropped, resized to 160x160 px. Model produces 512-float vector."),
    ("5", "L2 Normalise",      "Vector magnitude scaled to 1.0 — makes cosine similarity equal to dot product."),
]

# 5 pipeline rows: 0.74" each → 5×0.74 = 3.70, starting at 1.2, ending 4.90
# Cloud section: 4.10 to 6.95 (1.85"). Footer at 7.18. Total fine.
cy = 1.2
for num, title, detail in pipeline:
    y = cy
    box(s, 0.45, y, 12.43, 0.72, fill=CARD, border=BORDER)
    box(s, 0.45, y, 0.48, 0.72, fill=ACCENT)
    txt(s, num, 0.45, y + 0.17, 0.48, 0.38, size=15, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, title, 1.05, y + 0.08, 3.0, 0.32, size=14, bold=True, color=WHITE)
    txt(s, detail, 4.15, y + 0.08, 8.6, 0.56, size=13, color=MUTED)
    cy += 0.78

# Cloud section starts at cy + 0.18
cloud_y = cy + 0.18
box(s, 0.45, cloud_y, 12.43, 2.0, fill=DARK, border=BORDER)
txt(s, "Cloud Verification (Vercel / Neon PostgreSQL)",
    0.65, cloud_y + 0.10, 9.0, 0.34, size=14, bold=True, color=ACCENT)
cloud = [
    "1.  Lookup driver record by pinHash  ->  fetch encrypted embeddings_enc from Neon.",
    "2.  Decrypt with FLEET_EMBED_KEY  (AES-256-GCM — stored only in Vercel environment variables).",
    "3.  Cosine similarity against each of the 5 enrollment captures:  score = dot(liveVec, storedVec).",
    "4.  Return best (max) score.  Threshold:  score >= 0.75  =  MATCH   |   score < 0.75  =  FAIL.",
]
cy2 = cloud_y + 0.52
for line in cloud:
    txt(s, line, 0.65, cy2, 12.1, 0.32, size=12, color=BODY)
    cy2 += 0.34
footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Security Architecture
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Security Architecture", "Solution B — PIN + Face")

pillars = [
    ("PIN Protection",
     ["SHA-256(pin + pepper) — PIN never stored in plain text.",
      "Pepper value never enters the database; hash is one-way.",
      "6-digit PIN: 1,000,000 possible values with lockout enforcement."]),
    ("Biometric Privacy",
     ["Face images never leave the device — only math vectors (512 floats) transmitted.",
      "Embeddings at rest encrypted with AES-256-GCM using Vercel env key.",
      "No raw biometric data is ever stored or transmitted."]),
    ("Brute-Force Defence",
     ["Tablet lockout:  5 wrong PINs  ->  30-minute lock (SharedPreferences).",
      "Driver lockout:  5 wrong face attempts  ->  60-minute per-driver lock.",
      "Both locks use timestamp-based expiry — no server round-trip required."]),
    ("Key Management",
     ["FLEET_EMBED_KEY: 256-bit hex — stored only in Vercel environment variables.",
      "Never committed to source code or stored in the database.",
      "Key is rotatable without requiring drivers to re-enrol."]),
]

CW = 5.92
CH_BASE = 1.85  # fixed card height
for i, (title, lines) in enumerate(pillars):
    col = i % 2
    row = i // 2
    x = 0.45 + col * (CW + 0.54)
    y = 1.2  + row * (CH_BASE + 0.28)
    box(s, x, y, CW, CH_BASE, fill=CARD, border=BORDER)
    box(s, x, y, CW, 0.44, fill=ACCENT if col == 0 else MUTED)
    txt(s, title, x + 0.14, y + 0.06, CW - 0.28, 0.34, size=14, bold=True, color=BG)
    for j, line in enumerate(lines):
        txt(s, line, x + 0.14, y + 0.54 + j * 0.36, CW - 0.28, 0.32, size=12, color=BODY)

footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Comparison Table
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Head-to-Head Comparison")

rows = [
    ("Criterion",            "Passkey (WebAuthn)",    "NFC + PIN",            "PIN + Face"),
    ("Works on any tablet",  "No — device-bound",     "Yes — card-based",     "Yes — cloud lookup"),
    ("No special hardware",  "No — FIDO2 required",   "Yes — NFC reader",     "Yes — front camera"),
    ("Fleet enrollment",     "No — not supported",    "Yes — admin issues",   "Yes — self-register"),
    ("Two-factor auth",      "Yes — platform + PIN",  "Yes — have + know",    "Yes — know + are"),
    ("Liveness check",       "Platform-dependent",    "None",                 "Yes — blink detect"),
    ("Biometrics encrypted", "N/A",                   "N/A",                  "Yes — AES-256-GCM"),
    ("Risk if compromised",  "Low — device-bound",    "High — card loss",     "Low — nothing to lose"),
    ("Implementation",       "High complexity",       "Low",                  "Medium"),
]

COL_W = [3.55, 2.82, 2.60, 3.06]
COL_X = [0.45, 4.05, 6.92, 9.57]
ROW_H = 0.50
SY    = 1.18

for r, row in enumerate(rows):
    for c in range(4):
        x, w = COL_X[c], COL_W[c]
        y = SY + r * ROW_H
        if r == 0:
            bg = ACCENT if c > 0 else RGBColor(0x1e, 0x2d, 0x40)
        elif c == 3:
            bg = RGBColor(0x0e, 0x22, 0x2e) if r % 2 == 0 else RGBColor(0x11, 0x28, 0x36)
        else:
            bg = CARD if r % 2 == 0 else CARD_ALT
        box(s, x, y, w, ROW_H, fill=bg, border=BORDER)
        fc = BG if r == 0 and c > 0 else (ACCENT if c == 3 and r > 0 else BODY)
        txt(s, row[c], x + 0.1, y + 0.1, w - 0.15, 0.34,
            size=12, bold=(r == 0), color=fc)

footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Recommendation
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Recommendation")

box(s, 0.45, 1.18, 12.43, 0.60, fill=GREEN_BG, border=GREEN_FG)
txt(s, "Proceed with Solution B:  PIN + Face Verification",
    0.65, 1.25, 12.1, 0.48, size=21, bold=True, color=GREEN_FG)

reasons = [
    ("01", "No Hardware Dependency",
     ["Only requires the front camera already on fleet tablets — no FIDO2 chips or NFC cards."]),
    ("02", "True Two-Factor — Unforgeable",
     ["PIN (something you know) + face (something you are) — cannot be shared or delegated."]),
    ("03", "Fleet-Wide, Tablet-Agnostic",
     ["Embeddings in cloud (encrypted) — any tablet can verify any driver, no pairing needed."]),
    ("04", "Production-Ready Security Baseline",
     ["AES-256-GCM embeddings, SHA-256 PIN hashing, dual lockout, blink liveness detection."]),
    ("05", "NFC as Optional Enhancement",
     ["Solution A is fully built — can be added later as express re-auth shortcut for drivers."]),
]

# 5 cards: available height = 7.18 - 1.95 = 5.23". Each card = 0.88", gap = 0.10"
cy = 1.95
for num, title, lines in reasons:
    cy = numbered_card(s, num, title, lines, 0.45, cy, 12.43, card_h=0.88) + 0.10

footer(s)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Next Steps
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
slide_header(s, "Next Steps & Roadmap")

phases = [
    ("Phase 1 — Harden", "1 week",
     ["Move PIN pepper from source code to Vercel environment variable.",
      "Add per-driver salt to PIN hash schema in Neon PostgreSQL.",
      "Enable HTTPS certificate pinning inside FleetApiClient."]),
    ("Phase 2 — Validate", "2 weeks",
     ["Pilot with 10 drivers across 3 tablets in a real depot.",
      "Tune cosine similarity threshold (currently 0.75) using real driver dataset.",
      "Measure end-to-end auth time — target under 8 seconds."]),
    ("Phase 3 — Production", "4 weeks",
     ["Driver self-enrollment onboarding session (5-capture face enrol + 6-digit PIN).",
      "Admin dashboard for driver management and lockout override.",
      "Phased rollout to full fleet — depot by depot."]),
    ("Optional — NFC Fast Lane", "Parallel",
     ["Introduce NFC card as express re-auth shortcut for returning drivers.",
      "Fix PIN hashing and add card revocation table in NFC branch first."]),
]

CW = 5.92
CH = 2.0
for i, (phase, dur, lines) in enumerate(phases):
    col = i % 2
    row = i // 2
    x = 0.45 + col * (CW + 0.54)
    y = 1.2  + row * (CH + 0.30)
    box(s, x, y, CW, CH, fill=CARD, border=BORDER)
    box(s, x, y, CW, 0.46, fill=DARK, border=BORDER)
    txt(s, phase, x + 0.14, y + 0.08, CW - 1.1, 0.32, size=13, bold=True,
        color=ACCENT if i < 3 else MUTED)
    txt(s, dur, x + CW - 1.05, y + 0.08, 0.95, 0.30, size=11,
        color=MUTED, italic=True, align=PP_ALIGN.RIGHT)
    for j, line in enumerate(lines):
        txt(s, "- " + line, x + 0.14, y + 0.56 + j * 0.44, CW - 0.28, 0.38, size=12, color=BODY)

footer(s)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "WasteHero_Driver_Auth_POC.pptx")
prs.save(out)
print("Saved ->", out)
